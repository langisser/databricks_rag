#!/usr/bin/env python3
"""
EXTRACT ALL VOLUME FILES - Complete Document Processing
- Load and extract ALL files from Unity Catalog volume
- Support multiple file formats (DOCX, XLSX, PDF, TXT, etc.)
- Advanced chunking with semantic boundaries
- Create production-ready vector search index
- Comprehensive error handling and progress tracking

MANUAL EXECUTION:
python extract_all_volume_files.py
"""

import sys
import os
sys.path.append(os.path.join(os.path.dirname(__file__), '..', 'databricks_helper', 'helper_function'))

from databricks.connect import DatabricksSession
from databricks.vector_search.client import VectorSearchClient
import json
import time
import tempfile
import requests
import subprocess
from datetime import datetime

def install_required_libraries():
    """Install required libraries for document processing"""
    print("Installing required libraries...")
    libraries = [
        "python-docx",
        "pandas",
        "openpyxl",
        "xlrd",
        "PyPDF2",
        "pdfplumber",
        "python-pptx"
    ]

    for lib in libraries:
        try:
            print(f"  Installing {lib}...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", lib, "--quiet"])
        except subprocess.CalledProcessError as e:
            print(f"  WARNING: Failed to install {lib}: {e}")

    print("Library installation completed")

def download_file_from_volume(spark, volume_path, local_temp_path):
    """Download file from volume to local temp for processing"""
    try:
        binary_df = spark.read.format("binaryFile").load(volume_path)
        file_data = binary_df.collect()[0]

        with open(local_temp_path, 'wb') as f:
            f.write(file_data['content'])

        return True
    except Exception as e:
        print(f"    Download failed: {e}")
        return False

def extract_docx_content(local_file_path):
    """Extract full text from DOCX file"""
    try:
        from docx import Document

        doc = Document(local_file_path)
        full_text = []

        # Extract all paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                full_text.append(paragraph.text.strip())

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        full_text.append(cell.text.strip())

        content = "\\n".join(full_text)
        return content, "docx_full_extraction"

    except Exception as e:
        return f"[DOCX EXTRACTION ERROR: {str(e)}]", "docx_failed"

def extract_excel_content(local_file_path):
    """Extract full content from Excel file"""
    try:
        import pandas as pd

        all_content = []

        # Use context manager to ensure proper file closure
        with pd.ExcelFile(local_file_path) as excel_file:
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)

                # Add sheet header
                all_content.append(f"Sheet: {sheet_name}")

                # Add column headers
                headers = " | ".join([str(col) for col in df.columns])
                all_content.append(f"Columns: {headers}")

                # Add data rows
                for index, row in df.iterrows():
                    row_data = []
                    for value in row:
                        if pd.isna(value):
                            row_data.append("")
                        else:
                            row_data.append(str(value))

                    if any(cell.strip() for cell in row_data):
                        all_content.append(" | ".join(row_data))

                all_content.append("")  # Empty line between sheets

        content = "\\n".join(all_content)
        return content, "excel_full_extraction"

    except Exception as e:
        return f"[EXCEL EXTRACTION ERROR: {str(e)}]", "excel_failed"

def extract_pdf_content(local_file_path):
    """Extract text from PDF file"""
    try:
        import pdfplumber

        all_text = []
        with pdfplumber.open(local_file_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    all_text.append(f"[Page {page_num + 1}]")
                    all_text.append(text.strip())

        content = "\\n".join(all_text)
        return content, "pdf_pdfplumber_extraction"

    except Exception:
        # Fallback to PyPDF2
        try:
            import PyPDF2

            all_text = []
            with open(local_file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text:
                        all_text.append(f"[Page {page_num + 1}]")
                        all_text.append(text.strip())

            content = "\\n".join(all_text)
            return content, "pdf_pypdf2_extraction"

        except Exception as e:
            return f"[PDF EXTRACTION ERROR: {str(e)}]", "pdf_failed"

def extract_powerpoint_content(local_file_path):
    """Extract text from PowerPoint file"""
    try:
        from pptx import Presentation

        prs = Presentation(local_file_path)
        all_text = []

        for slide_num, slide in enumerate(prs.slides):
            all_text.append(f"[Slide {slide_num + 1}]")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    all_text.append(shape.text.strip())

        content = "\\n".join(all_text)
        return content, "pptx_extraction"

    except Exception as e:
        return f"[PPTX EXTRACTION ERROR: {str(e)}]", "pptx_failed"

def extract_text_file(local_file_path):
    """Extract content from text file"""
    try:
        with open(local_file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        return content, "text_file_utf8"
    except:
        try:
            with open(local_file_path, 'r', encoding='latin-1') as f:
                content = f.read()
            return content, "text_file_latin1"
        except Exception as e:
            return f"[TEXT EXTRACTION ERROR: {str(e)}]", "text_failed"

def create_semantic_chunks(content, file_name, word_count):
    """Create semantic chunks with advanced strategies"""
    import re

    # Determine chunk strategy based on content size
    if word_count < 500:
        chunk_size, overlap = 300, 50
    elif word_count < 2000:
        chunk_size, overlap = 500, 100
    elif word_count < 5000:
        chunk_size, overlap = 800, 150
    else:
        chunk_size, overlap = 1000, 200

    print(f"    Chunking strategy: {chunk_size} words, {overlap} overlap")

    # Clean content
    clean_content = re.sub(r'\\n+', ' ', content)
    clean_content = re.sub(r'\\s+', ' ', clean_content).strip()

    # Split into sentences
    sentences = re.split(r'[.!?]+\\s+', clean_content)
    sentences = [s.strip() for s in sentences if s.strip()]

    # Build chunks
    chunks = []
    current_chunk = ""
    current_words = 0
    chunk_index = 0

    for sentence in sentences:
        sentence_words = len(sentence.split())

        if current_words + sentence_words > chunk_size and current_chunk:
            chunks.append({
                "text": current_chunk.strip(),
                "word_count": current_words,
                "index": chunk_index
            })

            # Start new chunk with overlap
            if overlap > 0 and current_words > overlap:
                overlap_words = current_chunk.split()[-overlap:]
                current_chunk = " ".join(overlap_words) + " " + sentence
                current_words = len(overlap_words) + sentence_words
            else:
                current_chunk = sentence
                current_words = sentence_words

            chunk_index += 1
        else:
            if current_chunk:
                current_chunk += " " + sentence
            else:
                current_chunk = sentence
            current_words += sentence_words

    # Add final chunk
    if current_chunk:
        chunks.append({
            "text": current_chunk.strip(),
            "word_count": current_words,
            "index": chunk_index
        })

    return chunks, f"{chunk_size}w_{overlap}o"

def main():
    print("=" * 80)
    print("EXTRACT ALL VOLUME FILES - COMPLETE DOCUMENT PROCESSING")
    print("=" * 80)
    print(f"Started at: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

    try:
        # Install required libraries
        install_required_libraries()

        # Load configuration
        config_path = os.path.join(os.path.dirname(__file__), '..', 'databricks_helper', 'databricks_config', 'config.json')
        with open(config_path, 'r') as f:
            config = json.load(f)

        spark = DatabricksSession.builder.remote(
            host=config['databricks']['host'],
            token=config['databricks']['token'],
            cluster_id=config['databricks']['cluster_id']
        ).getOrCreate()

        namespace = config['rag_config']['full_namespace']
        volume_path = "/Volumes/sandbox/rdt_knowledge/rdt_document"

        print(f"\\nConfiguration:")
        print(f"  Volume: {volume_path}")
        print(f"  Namespace: {namespace}")

        # Discover ALL files in volume
        print(f"\\nDiscovering files in volume...")
        files_result = spark.sql(f"LIST '{volume_path}'")
        files = files_result.collect()

        print(f"Found {len(files)} files to process:")
        for i, file_info in enumerate(files, 1):
            file_name = file_info['name']
            file_size_kb = round(file_info['size'] / 1024, 1)
            file_type = file_name.split('.')[-1].lower() if '.' in file_name else 'unknown'
            print(f"  {i}. {file_name} ({file_type}, {file_size_kb} KB)")

        # Recreate tables for ALL files
        print(f"\\nRecreating tables for all documents...")

        spark.sql(f"DROP TABLE IF EXISTS {namespace}.all_documents")
        spark.sql(f"DROP TABLE IF EXISTS {namespace}.all_document_chunks")
        spark.sql(f"DROP TABLE IF EXISTS {namespace}.all_vs_chunks")

        # Create comprehensive tables
        spark.sql(f"""
            CREATE TABLE {namespace}.all_documents (
                document_id STRING,
                file_name STRING,
                file_path STRING,
                file_type STRING,
                file_size_kb DOUBLE,
                full_content STRING,
                content_length INT,
                word_count INT,
                paragraph_count INT,
                extraction_method STRING,
                processing_status STRING,
                created_timestamp TIMESTAMP
            )
            USING DELTA
        """)

        spark.sql(f"""
            CREATE TABLE {namespace}.all_document_chunks (
                chunk_id STRING,
                document_id STRING,
                chunk_index INT,
                chunk_text STRING,
                chunk_word_count INT,
                chunk_char_count INT,
                chunk_strategy STRING,
                semantic_density DOUBLE,
                source_file STRING,
                extraction_method STRING,
                created_timestamp TIMESTAMP
            )
            USING DELTA
        """)

        # Vector search compatible table
        spark.sql(f"""
            CREATE TABLE {namespace}.all_vs_chunks (
                chunk_id STRING,
                document_id STRING,
                chunk_index INT,
                chunk_text STRING,
                chunk_word_count INT,
                chunk_char_count INT,
                chunk_strategy STRING,
                semantic_density DOUBLE,
                source_file STRING,
                extraction_method STRING,
                created_timestamp TIMESTAMP
            )
            USING DELTA
            TBLPROPERTIES (delta.enableChangeDataFeed = true)
        """)

        print("Tables created successfully")

        # Process ALL files
        processed_count = 0
        total_chunks = 0
        total_words = 0
        success_count = 0
        failed_files = []

        for file_info in files:
            file_name = file_info['name']
            file_path = file_info['path']
            file_size = file_info['size']
            file_type = file_name.split('.')[-1].lower() if '.' in file_name else 'unknown'

            print(f"\\nProcessing [{processed_count + 1}/{len(files)}]: {file_name}")
            print(f"  Type: {file_type}, Size: {round(file_size/1024, 1)} KB")

            doc_id = f"all_doc_{int(time.time())}_{processed_count + 1}"

            with tempfile.TemporaryDirectory() as temp_dir:
                local_file_path = os.path.join(temp_dir, file_name)

                # Download file
                print(f"  Downloading from volume...")
                if not download_file_from_volume(spark, file_path, local_file_path):
                    print(f"  SKIP: Could not download {file_name}")
                    failed_files.append(f"{file_name} (download failed)")
                    processed_count += 1
                    continue

                # Extract content based on file type
                content = ""
                extraction_method = ""

                if file_type == 'docx':
                    print(f"  Extracting DOCX content...")
                    content, extraction_method = extract_docx_content(local_file_path)

                elif file_type in ['xlsx', 'xls']:
                    print(f"  Extracting Excel content...")
                    content, extraction_method = extract_excel_content(local_file_path)

                elif file_type == 'pdf':
                    print(f"  Extracting PDF content...")
                    content, extraction_method = extract_pdf_content(local_file_path)

                elif file_type in ['pptx', 'ppt']:
                    print(f"  Extracting PowerPoint content...")
                    content, extraction_method = extract_powerpoint_content(local_file_path)

                elif file_type in ['txt', 'md', 'csv']:
                    print(f"  Extracting text content...")
                    content, extraction_method = extract_text_file(local_file_path)

                else:
                    # Try as text file
                    print(f"  Attempting text extraction...")
                    content, extraction_method = extract_text_file(local_file_path)

                # Calculate content metrics
                if content and not content.startswith('['):  # Successful extraction
                    word_count = len(content.split())
                    char_count = len(content)
                    paragraph_count = len([p for p in content.split('\\n') if p.strip()])
                    processing_status = "success"

                    print(f"  SUCCESS: {word_count} words, {char_count} characters")
                    print(f"  Paragraphs: {paragraph_count}")

                    total_words += word_count
                    success_count += 1
                else:
                    word_count = 0
                    char_count = len(content) if content else 0
                    paragraph_count = 0
                    processing_status = "failed"
                    failed_files.append(f"{file_name} ({extraction_method})")
                    print(f"  FAILED: {content[:100]}...")

                # Insert document
                safe_content = content.replace("'", "''") if content else ""
                spark.sql(f"""
                    INSERT INTO {namespace}.all_documents VALUES (
                        '{doc_id}',
                        '{file_name}',
                        '{file_path}',
                        '{file_type}',
                        {file_size / 1024:.2f},
                        '{safe_content}',
                        {char_count},
                        {word_count},
                        {paragraph_count},
                        '{extraction_method}',
                        '{processing_status}',
                        current_timestamp()
                    )
                """)

                # Create chunks if successful
                if processing_status == "success" and word_count > 0:
                    print(f"  Creating semantic chunks...")

                    chunks, chunk_strategy = create_semantic_chunks(content, file_name, word_count)

                    for chunk in chunks:
                        chunk_id = f"{doc_id}_chunk_{chunk['index']:03d}"
                        chunk_text = chunk['text']
                        chunk_words = chunk['word_count']
                        chunk_chars = len(chunk_text)

                        # Calculate semantic density
                        words = chunk_text.lower().split()
                        unique_words = len(set(words))
                        semantic_density = unique_words / len(words) if words else 0

                        safe_chunk_text = chunk_text.replace("'", "''")

                        # Insert into chunks table
                        spark.sql(f"""
                            INSERT INTO {namespace}.all_document_chunks VALUES (
                                '{chunk_id}',
                                '{doc_id}',
                                {chunk['index']},
                                '{safe_chunk_text}',
                                {chunk_words},
                                {chunk_chars},
                                '{chunk_strategy}',
                                {semantic_density:.4f},
                                '{file_name}',
                                '{extraction_method}',
                                current_timestamp()
                            )
                        """)

                        # Insert into vector search compatible table
                        spark.sql(f"""
                            INSERT INTO {namespace}.all_vs_chunks VALUES (
                                '{chunk_id}',
                                '{doc_id}',
                                {chunk['index']},
                                '{safe_chunk_text}',
                                {chunk_words},
                                {chunk_chars},
                                '{chunk_strategy}',
                                {semantic_density:.4f},
                                '{file_name}',
                                '{extraction_method}',
                                current_timestamp()
                            )
                        """)

                    total_chunks += len(chunks)
                    print(f"  Created {len(chunks)} chunks")

                processed_count += 1

        spark.stop()

        # Create vector search index
        print(f"\\n" + "=" * 80)
        print("CREATING VECTOR SEARCH INDEX WITH ALL DOCUMENTS")
        print("=" * 80)

        # Vector Search client
        vsc = VectorSearchClient(
            workspace_url=config['databricks']['host'],
            personal_access_token=config['databricks']['token'],
            disable_notice=True
        )

        # Create comprehensive index
        all_index_name = f"{namespace}.all_content_rag_index"
        source_table = f"{namespace}.all_vs_chunks"
        endpoint_name = config['rag_config']['vector_search_endpoint']
        embedding_model = config['rag_config']['embedding_model']

        try:
            # Delete existing index if exists
            try:
                vsc.delete_index(all_index_name)
                print("Deleted existing index")
                time.sleep(10)
            except:
                print("No existing index to delete")

            # Create new comprehensive index
            print(f"Creating index: {all_index_name}")

            index = vsc.create_delta_sync_index(
                endpoint_name=endpoint_name,
                index_name=all_index_name,
                source_table_name=source_table,
                pipeline_type="TRIGGERED",
                primary_key="chunk_id",
                embedding_source_column="chunk_text",
                embedding_model_endpoint_name=embedding_model
            )

            print("SUCCESS: Vector index creation initiated")

        except Exception as e:
            print(f"WARNING: Vector index creation failed: {e}")

        # Update configuration
        config['rag_config']['all_content_index'] = all_index_name
        config['rag_config']['all_content_table'] = source_table

        with open(config_path, 'w') as f:
            json.dump(config, f, indent=2)

        # Final summary
        print(f"\\n" + "=" * 80)
        print("ALL VOLUME FILES EXTRACTION COMPLETED")
        print("=" * 80)
        print(f"Completed at: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

        print(f"\\nPROCESSING RESULTS:")
        print(f"  Total files found: {len(files)}")
        print(f"  Successfully processed: {success_count}")
        print(f"  Failed processing: {len(failed_files)}")
        print(f"  Total words extracted: {total_words:,}")
        print(f"  Total chunks created: {total_chunks}")

        if failed_files:
            print(f"\\nFAILED FILES:")
            for failed_file in failed_files:
                print(f"  - {failed_file}")

        print(f"\\nTABLES CREATED:")
        print(f"  Documents: {namespace}.all_documents")
        print(f"  Chunks: {namespace}.all_document_chunks")
        print(f"  Vector Search: {namespace}.all_vs_chunks")

        print(f"\\nVECTOR SEARCH INDEX:")
        print(f"  Index: {all_index_name}")
        print(f"  Status: Building (wait 15-20 minutes)")

        print(f"\\nREADY FOR: Production RAG with complete document collection!")

        return True

    except Exception as e:
        print(f"\\nERROR: All files extraction failed: {e}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)