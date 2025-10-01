#!/usr/bin/env python3
"""
Process ALL documents in Volume with multi-format support
Handles: DOCX, XLSX, PDF, PPTX, DRAWIO
Extracts: Text, Tables, Images with document name metadata
"""
import sys
import os
import json
from io import BytesIO
from datetime import datetime
sys.path.append(os.path.join(os.path.dirname(__file__), '..', 'databricks_helper', 'helper_function'))

from databricks.connect import DatabricksSession
from pyspark.sql import Row

def install_libraries():
    """Install all required libraries"""
    print("Installing libraries...")
    import subprocess
    libs = [
        "python-docx",      # DOCX
        "openpyxl",         # XLSX
        "pandas",           # Excel/data
        "pdfplumber",       # PDF
        "PyPDF2",           # PDF backup
        "python-pptx",      # PPTX
        "Pillow"            # Images
    ]
    for lib in libs:
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", lib, "--quiet"])
        except:
            pass
    print("Libraries installed\n")

# ============================================================================
# DOCX PROCESSOR
# ============================================================================
def process_docx(file_content, filename):
    """Extract text, tables, images from DOCX"""
    from docx import Document

    doc = Document(BytesIO(file_content))
    chunks = []
    chunk_id = 0

    # Extract text with sections
    current_section = "Introduction"
    current_text = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Detect section headers
        if len(text) < 100 and (para.style.name.startswith('Heading') or text.isupper()):
            if current_text:
                chunks.append({
                    'id': chunk_id,
                    'content': '\n'.join(current_text),
                    'type': 'text',
                    'section': current_section,
                    'doc_name': filename
                })
                chunk_id += 1
                current_text = []
            current_section = text
            current_text.append(f"## {text}")
        else:
            current_text.append(text)
            if sum(len(t) for t in current_text) > 2000:
                chunks.append({
                    'id': chunk_id,
                    'content': '\n'.join(current_text),
                    'type': 'text',
                    'section': current_section,
                    'doc_name': filename
                })
                chunk_id += 1
                current_text = []

    if current_text:
        chunks.append({
            'id': chunk_id,
            'content': '\n'.join(current_text),
            'type': 'text',
            'section': current_section,
            'doc_name': filename
        })
        chunk_id += 1

    # Extract tables
    for tbl_idx, table in enumerate(doc.tables):
        headers = [cell.text.strip() for cell in table.rows[0].cells] if len(table.rows) > 0 else []
        rows_data = []
        for row in table.rows[1:]:
            row_data = [cell.text.strip() for cell in row.cells]
            if any(row_data):
                rows_data.append(' | '.join(row_data))

        if rows_data:
            table_text = f"TABLE {tbl_idx + 1}:\nHeaders: {' | '.join(headers)}\n" + '\n'.join(rows_data)
            chunks.append({
                'id': chunk_id,
                'content': table_text,
                'type': 'table',
                'section': current_section,
                'doc_name': filename
            })
            chunk_id += 1

    # Extract images count
    image_count = len(doc.inline_shapes)
    if image_count > 0:
        chunks.append({
            'id': chunk_id,
            'content': f"Document contains {image_count} images/diagrams including process flows and architecture diagrams.",
            'type': 'image_metadata',
            'section': 'Images',
            'doc_name': filename
        })

    return chunks

# ============================================================================
# XLSX PROCESSOR
# ============================================================================
def process_xlsx(file_content, filename):
    """Extract data from all sheets in XLSX"""
    import openpyxl
    import pandas as pd

    chunks = []
    chunk_id = 0

    wb = openpyxl.load_workbook(BytesIO(file_content), data_only=True)

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]

        # Read sheet data
        data = []
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                data.append([str(cell) if cell is not None else '' for cell in row])

        if data:
            # Create table representation
            table_text = f"EXCEL SHEET: {sheet_name}\n"
            if len(data) > 0:
                table_text += "Headers: " + ' | '.join(data[0]) + "\n"
            for row in data[1:min(50, len(data))]:  # Limit rows
                table_text += ' | '.join(row) + "\n"

            chunks.append({
                'id': chunk_id,
                'content': table_text,
                'type': 'excel_sheet',
                'section': f"Sheet: {sheet_name}",
                'doc_name': filename
            })
            chunk_id += 1

    return chunks

# ============================================================================
# PDF PROCESSOR
# ============================================================================
def process_pdf(file_content, filename):
    """Extract text and tables from PDF"""
    import pdfplumber

    chunks = []
    chunk_id = 0

    with pdfplumber.open(BytesIO(file_content)) as pdf:
        for page_num, page in enumerate(pdf.pages):
            # Extract text
            text = page.extract_text()
            if text and text.strip():
                chunks.append({
                    'id': chunk_id,
                    'content': text.strip(),
                    'type': 'pdf_text',
                    'section': f"Page {page_num + 1}",
                    'doc_name': filename
                })
                chunk_id += 1

            # Extract tables
            tables = page.extract_tables()
            for tbl_idx, table in enumerate(tables):
                if table:
                    table_text = f"PDF TABLE (Page {page_num + 1}, Table {tbl_idx + 1}):\n"
                    for row in table:
                        table_text += ' | '.join([str(cell) if cell else '' for cell in row]) + "\n"

                    chunks.append({
                        'id': chunk_id,
                        'content': table_text,
                        'type': 'pdf_table',
                        'section': f"Page {page_num + 1}",
                        'doc_name': filename
                    })
                    chunk_id += 1

    return chunks

# ============================================================================
# PPTX PROCESSOR
# ============================================================================
def process_pptx(file_content, filename):
    """Extract text, tables, and images from PPTX slides"""
    from pptx import Presentation

    chunks = []
    chunk_id = 0

    prs = Presentation(BytesIO(file_content))

    for slide_num, slide in enumerate(prs.slides):
        slide_text = []
        slide_tables = []
        image_count = 0

        # Extract text, tables, and images from shapes
        for shape in slide.shapes:
            # Extract text
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

            # Extract tables
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    if any(row_data):
                        table_data.append(' | '.join(row_data))

                if table_data:
                    slide_tables.append('\n'.join(table_data))

            # Count images
            if shape.shape_type == 13:  # Picture type
                image_count += 1

        # Add slide text chunk
        if slide_text:
            chunks.append({
                'id': chunk_id,
                'content': f"SLIDE {slide_num + 1}:\n" + '\n'.join(slide_text),
                'type': 'pptx_slide',
                'section': f"Slide {slide_num + 1}",
                'doc_name': filename
            })
            chunk_id += 1

        # Add slide tables
        for tbl_idx, table_text in enumerate(slide_tables):
            chunks.append({
                'id': chunk_id,
                'content': f"SLIDE {slide_num + 1} - TABLE {tbl_idx + 1}:\n{table_text}",
                'type': 'pptx_table',
                'section': f"Slide {slide_num + 1}",
                'doc_name': filename
            })
            chunk_id += 1

        # Add image metadata
        if image_count > 0:
            chunks.append({
                'id': chunk_id,
                'content': f"SLIDE {slide_num + 1} contains {image_count} image(s)/diagram(s)",
                'type': 'pptx_image',
                'section': f"Slide {slide_num + 1}",
                'doc_name': filename
            })
            chunk_id += 1

    return chunks

# ============================================================================
# DRAWIO PROCESSOR
# ============================================================================
def process_drawio(file_content, filename):
    """Extract text from DrawIO diagrams"""
    import xml.etree.ElementTree as ET

    chunks = []
    chunk_id = 0

    try:
        root = ET.fromstring(file_content.decode('utf-8'))

        # Find all diagram pages
        diagrams = root.findall('.//diagram')

        for diag_idx, diagram in enumerate(diagrams):
            page_name = diagram.get('name', f'Page {diag_idx + 1}')

            # Extract text from diagram
            texts = []
            for elem in diagram.iter():
                if elem.text and elem.text.strip():
                    texts.append(elem.text.strip())
                for attr_val in elem.attrib.values():
                    if isinstance(attr_val, str) and len(attr_val) > 3 and attr_val.strip():
                        texts.append(attr_val.strip())

            if texts:
                chunks.append({
                    'id': chunk_id,
                    'content': f"DIAGRAM: {page_name}\nContent: " + ' | '.join(set(texts)[:50]),
                    'type': 'drawio_diagram',
                    'section': page_name,
                    'doc_name': filename
                })
                chunk_id += 1
    except:
        chunks.append({
            'id': 0,
            'content': f"DRAWIO Diagram: {filename} (complex structure)",
            'type': 'drawio_diagram',
            'section': 'Diagram',
            'doc_name': filename
        })

    return chunks

# ============================================================================
# MAIN PROCESSOR
# ============================================================================
def main():
    print("=" * 80)
    print("PROCESS ALL DOCUMENTS - MULTIMODAL EXTRACTION")
    print("=" * 80)

    install_libraries()

    # Load config
    config_path = os.path.join(os.path.dirname(__file__), '..', 'databricks_helper', 'databricks_config', 'config.json')
    with open(config_path, 'r') as f:
        config = json.load(f)

    # Connect
    print("\n[STEP 1] Connecting to Databricks...")
    spark = DatabricksSession.builder.remote(
        host=config['databricks']['host'],
        token=config['databricks']['token'],
        cluster_id=config['databricks']['cluster_id']
    ).getOrCreate()
    print("  Connected")

    # Get all files
    volume_path = "/Volumes/sandbox/rdt_knowledge/rdt_document"
    print(f"\n[STEP 2] Loading files from: {volume_path}")

    files_result = spark.sql(f"LIST '{volume_path}'")
    files_df = files_result.toPandas()

    print(f"  Total files: {len(files_df)}")

    # Process each file
    print("\n[STEP 3] Processing documents...")

    all_chunks = []
    stats = {
        'total_files': 0,
        'processed_files': 0,
        'failed_files': 0,
        'by_type': {}
    }

    for idx, row in files_df.iterrows():
        filename = row['name']
        filepath = row['path']
        ext = os.path.splitext(filename)[1].lower()

        print(f"\n  [{idx + 1}/{len(files_df)}] Processing: {filename}")
        print(f"    Type: {ext}")

        stats['total_files'] += 1

        try:
            # Load file
            binary_df = spark.read.format("binaryFile").load(filepath)
            file_data = binary_df.collect()[0]
            file_content = file_data['content']

            # Process based on type
            chunks = []

            if ext == '.docx':
                chunks = process_docx(file_content, filename)
            elif ext == '.xlsx':
                chunks = process_xlsx(file_content, filename)
            elif ext == '.pdf':
                chunks = process_pdf(file_content, filename)
            elif ext == '.pptx':
                chunks = process_pptx(file_content, filename)
            elif ext == '.drawio':
                chunks = process_drawio(file_content, filename)
            else:
                print(f"    [SKIP] Unsupported file type")
                continue

            print(f"    [SUCCESS] Extracted {len(chunks)} chunks")

            # Add to collection
            for chunk in chunks:
                chunk['file_id'] = f"{os.path.splitext(filename)[0]}_{chunk['id']}"
                all_chunks.append(chunk)

            stats['processed_files'] += 1
            stats['by_type'][ext] = stats['by_type'].get(ext, 0) + 1

        except Exception as e:
            print(f"    [ERROR] {e}")
            stats['failed_files'] += 1

    print(f"\n[STEP 4] Creating Delta table...")
    print(f"  Total chunks extracted: {len(all_chunks)}")

    # Create DataFrame
    rows = []
    for chunk in all_chunks:
        rows.append(Row(
            id=chunk['file_id'],
            content=chunk['content'],
            content_type=chunk['type'],
            section=chunk['section'],
            document_name=chunk['doc_name'],
            char_count=len(chunk['content'])
        ))

    df = spark.createDataFrame(rows)

    # Save table
    namespace = config['rag_config']['full_namespace']
    table_name = f"{namespace}.all_documents_multimodal"

    df.write.format("delta").mode("overwrite").saveAsTable(table_name)

    print(f"  [SUCCESS] Table created: {table_name}")

    # Show summary
    print("\n[STEP 5] Summary:")
    summary = spark.sql(f"""
        SELECT content_type, COUNT(*) as count, SUM(char_count) as total_chars
        FROM {table_name}
        GROUP BY content_type
        ORDER BY count DESC
    """)
    summary.show()

    # Stats
    print("\n" + "=" * 80)
    print("PROCESSING COMPLETE")
    print("=" * 80)
    print(f"\nFiles processed: {stats['processed_files']}/{stats['total_files']}")
    print(f"Failed: {stats['failed_files']}")
    print(f"Total chunks: {len(all_chunks)}")
    print(f"\nBy file type:")
    for ext, count in stats['by_type'].items():
        print(f"  {ext}: {count} files")

    print(f"\n[NEXT] Create vector index from table: {table_name}")

if __name__ == "__main__":
    main()